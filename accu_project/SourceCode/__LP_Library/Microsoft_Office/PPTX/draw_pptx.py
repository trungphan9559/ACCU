import re

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN,MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE,XL_LEGEND_POSITION,XL_LABEL_POSITION,XL_CATEGORY_TYPE,XL_TICK_MARK
from pprint import pprint


class DrawPPTX:
	__slide = ''

	def __init__(self,slide):
		self.__slide = slide

	#CHART
	#Dùng để vẽ line,bar,column,column_stack chart
	def __create_normal_chart(self,chart_type,categories,dict_data,left,top,width,height,
														legend_font=12,title=None,has_legend=True,has_data_labels=False,number_format=''):

		#Nếu chart không có categories thì bỏ qua không vẽ
		if len(categories) == 0:
			return

		#1)CATEGORIES
		chart_data = CategoryChartData()
		chart_data.categories = categories

		#2)DATA IN CHART
		for key,value in dict_data.items():
			chart_data.add_series(key, value,number_format=number_format)

		#3)Thêm chart vào slide
		left = Inches(left)
		top = Inches(top)
		width = Inches(width)
		height = Inches(height)
		chart = self.__slide.shapes.add_chart(
				chart_type, left, top, width, height, chart_data
		).chart

		#3.1)CHART TITLE
		if title != None:
			chart.has_title = True
			chart.chart_title.text_frame.text = title
			chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
			chart.chart_title.text_frame.paragraphs[0].font.name = "Yu Gothic"

		#3.2)Categories
		chart.has_legend = has_legend
		if has_legend:
			chart.legend.font.size = Pt(legend_font)
			chart.legend.font.name = "Yu Gothic"
			chart.legend.position = XL_LEGEND_POSITION.TOP
			chart.legend.include_in_layout = False

		#3.3)Data ở trên cái cột
		if has_data_labels:
			plot = chart.plots[0]
			plot.has_data_labels = True
			data_labels = plot.data_labels
			data_labels.font.size = Pt(8)
			data_labels.font.name = "Yu Gothic"
			data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)

			# data_labels.position = XL_LABEL_POSITION.ABOVE

		#3.4)Trục Y
		value_axis = chart.value_axis
		value_axis.minimum_scale = 0
		tick_labels = value_axis.tick_labels
		tick_labels.font.size = Pt(9)
		tick_labels.font.name = "Yu Gothic"

		#3.4.2)Number Format
		value = list(dict_data.values())[0][0]

		if type(value) == int and number_format == '':
			tick_labels.number_format = "#,##0"

		#3.4.3)Nếu giá trị nhỏ hơn 1 thì tạo maximum bằng 8
		maximum_scale = 8
		max_value = max(list(dict_data.values())[0])
		if type(max_value) == int and max_value <= maximum_scale:
			value_axis.maximum_scale = maximum_scale

		#3.5)Trục X
		category_axis = chart.category_axis
		# category_axis.has_major_gridlines = True
		# category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
		category_axis.tick_labels.font.size = Pt(9)
		category_axis.tick_labels.font.name = "Yu Gothic"

	#1)TABLE (label là cái nằm ngang, categories là cái ô vuông có màu màu)
	#1.1)Dùng cho pie chart hoặc bảng có nhiều cột
	def __create_chart_table_follow_categories(self,header,categories,dict_data,left,top,cell_width,cell_height,font_size=12,table_columns_width={}):
		#1)Bảng này chỉ có 2 cột thôi
		row = len(dict_data)
		col = 2
		data_table = [[0 for c in range(col)] for r in range(row)]

		#2)Đưa dữ liệu vào data_table
		row = 0
		for key,value in dict_data.items():
				data_table[row][0] = key
				for data in value:
					data_table[row][1] += data
				row += 1

		#3)Vẽ table
		self.create_table(header,data_table,left,top,cell_width,cell_height,font_size=font_size,table_columns_width=table_columns_width)
	
	#1.2)Dùng để cho bảng chỉ có 1 categories
	def __create_chart_table_follow_label(self,header,categories,dict_data,left,top,cell_width,cell_height,font_size=12,table_columns_width={},dict_data_table=None):
		#1)Bảng này cũng chỉ có 2 cột
		row = len(categories)
		col = 2
		data_table = [[0 for c in range(col)] for r in range(row)]

		#2)Đưa dữ liệu vào data_table
		for r in range(0,row):
			for list_value in dict_data.values():
				data_table[r][1] += list_value[r]
			data_table[r][0] = categories[r]

		
		#2.1) Custom thêm 	
		if dict_data_table != None:
			data_table_ = [[0 for c in range(col)] for r in range(row)]

			for r in range(0,row):
				for list_value in dict_data_table.values():
					data_table_[r][1] = list_value[r]
				data_table_[r][0] = categories[r]
			
			data_table = data_table_
			
			
		#3)Vẽ table
		self.create_table(header,data_table,left,top,cell_width,cell_height,font_size=font_size,table_columns_width=table_columns_width)
	
	#1.3)Dùng để vẽ theo cả label và categories
	def create_chart_table_follow_both_label_and_categories(self,header,categories,dict_data,left,top,cell_width,cell_height,font_size=12,table_columns_width={}):
		#1)Bảng này cũng chỉ có 2 cột
		row = len(dict_data)
		col = len(header)
		data_table = [[0 for c in range(col)] for r in range(row)]

		#2)Đưa dữ liệu vào data_table
		r = 0
		for key,value in dict_data.items():
			data_table[r][0] = key
			for c,data in enumerate(value):
				data_table[r][c + 1] = data
			r += 1

		#3)Vẽ table
		self.create_table(header,data_table,left,top,cell_width,cell_height,font_size=12,table_columns_width=table_columns_width)
	
	#1.4)Chỉ dùng để vẽ bảng cho pie chart
	def create_pie_chart_table(self,header,categories,dict_data,left,top,cell_width,cell_height,font_size=12,table_columns_width={}):
		#1)Tạo ra layout trước
		row = len(categories)
		col = len(header)
		data_table = [[0 for c in range(col)] for r in range(row)]

		#1.1)Tổng dữ liệu của cái bánh
		total_data = list(dict_data.values())
		total_data = total_data[0]
		total_data = sum(total_data)
		
		#2)Đưa dữ liệu vào data_table
		for r in range(0,row):
			for list_value in dict_data.values():
				data = list_value[r]

				data_table[r][1] = data

				rate = data/total_data*100 if total_data != 0 else 0
				rate = round(rate,2)
				rate = str(rate) + '%'
				data_table[r][2] = rate

			data_table[r][0] = categories[r]

		#3)Vẽ table
		self.create_table(header,data_table,left,top,cell_width,cell_height,font_size=font_size,table_columns_width=table_columns_width)

	#1.5)Dùng để cho bảng chỉ có 1 categories, dùng cho những table là đoạn text chứ không phải là số
	def create_table_follow_label(self,header,categories,dict_data,left,top,cell_width,cell_height,font_size=12,table_columns_width={}):
		#1)Bảng này cũng chỉ có 2 cột
		row = len(categories)
		col = 2
		data_table = [[0 for c in range(col)] for r in range(row)]

		#2)Đưa dữ liệu vào data_table
		for r in range(0,row):
			for list_value in dict_data.values():
				data_table[r][1] = list_value[r]
			data_table[r][0] = categories[r]

		#3)Vẽ table
		self.create_table(header,data_table,left,top,cell_width,cell_height,font_size=font_size,table_columns_width=table_columns_width)

	def create_table(self,header,data_table,left,top,cell_width,cell_height,font_size=12,alignment="RIGHT",alignment_col=0,table_columns_width={}):
		try:
			#1)Khởi tạo table và đặt vị trí
			row = len(data_table) + 1
			col = len(header)
			
			x, y, cx, cy = Inches(left), Inches(top), Inches(cell_width), Inches(cell_height)
			shape = self.__slide.shapes.add_table(row, col, x, y, cx, cy)
			table = shape.table

			#2)Đưa dữ liệu vào table
			#2.1)Header
			r = c = 0
			for data in header:
				cell = table.cell(r, c)
				cell.text = str(data)
				# cell.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

				p = cell.text_frame.paragraphs[0]
				# p.text = str(data)
				p.font.size = Pt(font_size)
				p.font.name = "Yu Gothic"
				p.alignment = PP_ALIGN.CENTER

				c += 1
			
			#2.2)Set độ rộng cho từng cột
			number_col_has_width = len(table_columns_width)

			sum_col_has_width = sum(table_columns_width.values())

			width_of_columns_unset = (cell_width - sum_col_has_width)/(col - number_col_has_width)

			for (index,column) in enumerate(table.columns):
				width = width_of_columns_unset

				if index in table_columns_width.keys():
					width = table_columns_width[index]

				column.width = Inches(width)

			#2.3)Data table
			r = 1
			c = 0
			regex = re.compile(r'[\n\r\t\b]')
			
			for tr in data_table:
				for td in tr:
					#2.2.1)Nếu td là số 123456 thì đổi nó thành 123,456
					if isinstance(td,int) or isinstance(td,float):
						td = format(td,',')
					#2.2.2)Xóa ký tự control character
					td = regex.sub("", td)
		
					cell = table.cell(r, c)
					cell.text = str(td)

					p = cell.text_frame.paragraphs[0]
					# p.text = str(td)
					p.font.size = Pt(font_size)
					p.font.name = "Yu Gothic"
					if alignment == "RIGHT" and c > alignment_col:
						p.alignment = PP_ALIGN.RIGHT
					if alignment == "CENTER" and c > alignment_col:
						p.alignment = PP_ALIGN.CENTER
					c += 1
				r += 1
				c = 0
		except Exception as e:
			print(e)

	#IMG
	def create_img(self,img_path,left,top,height):
		left = Inches(left)
		top = Inches(top)
		height = Inches(height)
		self.__slide.shapes.add_picture(img_path, left, top, height=height)
	
	#TEXT
	def create_one_line_text(self,text_1,left,top,width,height,back_ground_color=RGBColor(255, 255, 255),
													font_size=14,text_color=RGBColor(0, 0, 0),
													font_bold = False,font_name = "Yu Gothic",alignment=False,hyperlink=None):
		#1)TEXT BOX COLOR
		left = Inches(left)
		top = Inches(top)
		width = Inches(width)
		height = Inches(height)

		txBox = self.__slide.shapes.add_textbox(left, top, width, height)
		#1.1)FILL BACK GROUND COLOR
		txBox.fill.solid()
		txBox.fill.fore_color.rgb = back_ground_color

		#2)Tạo đoạn text bên trong
		tf = txBox.text_frame
		tf.word_wrap = True

		#2.1) Dòng text đầu tiên
		p = tf.paragraphs[0]
		p.text = str(text_1)
		p.font.size = Pt(font_size)
		p.font.color.rgb = text_color
		p.font.bold = font_bold
		p.font.name = font_name
		#=== Căn giữa ===
		if alignment == True:
			p.alignment = PP_ALIGN.CENTER
		#=== Tạo hyper link ===
		if hyperlink != None:
			r = p.add_run()
			r.text = hyperlink['text']
			r.hyperlink.address = hyperlink['address']

	def create_two_lines_text(self,text_1,text_2,left,top,width,height,back_ground_color=RGBColor(255, 255, 255),
														font_size_text_1=14,font_size_text_2=14,
														text_color=RGBColor(0, 0, 0),font_bold = False,
														font_name = "Yu Gothic",alignment=False):
		# 1)Vị trí của đoạn text
		left = Inches(left)
		top = Inches(top)
		width = Inches(width)
		height = Inches(height)

		txBox = self.__slide.shapes.add_textbox(left, top, width, height)
		# 1.1)Chỉnh màu background
		txBox.fill.solid()
		txBox.fill.fore_color.	rgb = back_ground_color

		# 2)Tạo các dòng text
		tf = txBox.text_frame
		tf.word_wrap = True
		# 2.1) Dòng 1 trong text_box
		p = tf.paragraphs[0]
		p.text = str(text_1)
		p.font.size = Pt(font_size_text_1)
		p.font.color.rgb = text_color
		p.font.bold = font_bold
		p.font.name = font_name
		#=== Căn giữa ===
		if alignment == True:
			p.alignment = PP_ALIGN.CENTER

		# 2.2) Dòng 2 trong text_box
		p = tf.add_paragraph()
		p.text = str(text_2)
		p.font.size = Pt(font_size_text_2)
		p.font.color.rgb = text_color
		p.font.bold = font_bold
		p.font.name = font_name
		#=== Căn giữa ===
		if alignment == True:
			p.alignment = PP_ALIGN.CENTER

	def create_three_lines_text(self,text_1,text_2,text_3,left,top,width,height,back_ground_color=RGBColor(255, 255, 255),
														font_size_text_1=14,font_size_text_2=14,font_size_text_3=14,
														text_color=RGBColor(0, 0, 0),font_bold = False,
														font_name = "Yu Gothic",alignment=False):
		# 1)Vị trí của đoạn text
		left = Inches(left)
		top = Inches(top)
		width = Inches(width)
		height = Inches(height)

		txBox = self.__slide.shapes.add_textbox(left, top, width, height)
		# 1.1)Chỉnh màu background
		txBox.fill.solid()
		txBox.fill.fore_color.	rgb = back_ground_color

		# 2)Tạo các dòng text
		tf = txBox.text_frame
		tf.word_wrap = True
		# 2.1) Dòng 1 trong text_box
		p = tf.paragraphs[0]
		p.text = str(text_1)
		p.font.size = Pt(font_size_text_1)
		p.font.color.rgb = text_color
		p.font.bold = font_bold
		p.font.name = font_name
		p.line_spacing  = 1.5
		#=== Căn giữa ===
		if alignment == True:
			p.alignment = PP_ALIGN.CENTER

		# 2.2) Dòng 2 trong text_box
		p = tf.add_paragraph()
		p.text = str(text_2)
		p.font.size = Pt(font_size_text_2)
		p.font.color.rgb = text_color
		p.font.bold = font_bold
		p.font.name = font_name
		p.line_spacing  = 1.5
		#=== Căn giữa ===
		if alignment == True:
			p.alignment = PP_ALIGN.CENTER

		# 2.2) Dòng 3 trong text_box
		p = tf.add_paragraph()
		p.text = str(text_3)
		p.font.size = Pt(font_size_text_3)
		p.font.color.rgb = text_color
		p.font.bold = font_bold
		p.font.name = font_name
		p.line_spacing  = 1.5
		#=== Căn giữa ===
		if alignment == True:
			p.alignment = PP_ALIGN.CENTER

	def create_multi_lines_text(self,left,top,width,height,list_text=[],
														back_ground_color=RGBColor(255, 255, 255),
														text_color=RGBColor(0, 0, 0),font_bold = False,
														font_name = "Yu Gothic",alignment=False):
		'''
		list_text=[{
			"text" : "ABC",
			"font_size" : 14,
			}]
		'''
		# 1)Vị trí của đoạn text
		left = Inches(left)
		top = Inches(top)
		width = Inches(width)
		height = Inches(height)

		txBox = self.__slide.shapes.add_textbox(left, top, width, height)
		# 1.1)Chỉnh màu background
		txBox.fill.solid()
		txBox.fill.fore_color.	rgb = back_ground_color

		# 2)Tạo các dòng text
		tf = txBox.text_frame
		tf.word_wrap = True

		# 2.1) Dòng 1 trong text_box
		for item in list_text:
			p = tf.add_paragraph()
			p.text = str(item['text'])
			p.font.size = Pt(item['font_size'])
			p.font.color.rgb = text_color
			p.font.bold = font_bold
			p.font.name = font_name
			p.line_spacing  = 1.5
			#=== Căn giữa ===
			if alignment == True:
				p.alignment = PP_ALIGN.CENTER

		# # 2.2) Dòng 2 trong text_box
		# p = tf.add_paragraph()
		# p.text = str(text_2)
		# p.font.size = Pt(font_size_text_2)
		# p.font.color.rgb = text_color
		# p.font.bold = font_bold
		# p.font.name = font_name
		# p.line_spacing  = 1.5
		# #=== Căn giữa ===
		# if alignment == True:
		# 	p.alignment = PP_ALIGN.CENTER

		# # 2.2) Dòng 3 trong text_box
		# p = tf.add_paragraph()
		# p.text = str(text_3)
		# p.font.size = Pt(font_size_text_3)
		# p.font.color.rgb = text_color
		# p.font.bold = font_bold
		# p.font.name = font_name
		# p.line_spacing  = 1.5
		# #=== Căn giữa ===
		# if alignment == True:
		# 	p.alignment = PP_ALIGN.CENTER
		
	'''
	Hàm vẽ chart cấp 1 chỉ có 1 chức năng là vẽ chart
	'''
	# PIE CHART
	def create_pie_chart(self,categories,dict_data,left,top,width,height,legend_font=12,title=None,has_legend=True,has_data_labels=False):
		#0)Tổng dữ liệu của cái bánh
		total_data = list(dict_data.values())
		total_data = total_data[0]
		total_data = sum(total_data)
		if total_data == 0:
			return

		#1)Thêm categori vào chart
		chart_data = CategoryChartData()
		chart_data.categories = categories
		#1.1)Thêm dữ liệu vào
		for key,value in dict_data.items():
			chart_data.add_series(key, [x/total_data for x in value])

		#2)Thêm chart vào slide
		left = Inches(left)
		top = Inches(top)
		width = Inches(width)
		height = Inches(height)
		chart = self.__slide.shapes.add_chart(
				XL_CHART_TYPE.PIE, left, top, width, height, chart_data
		).chart

		#2.1)CHART TITLE
		if title != None:
			chart.has_title = True
			chart.chart_title.text_frame.text = title
			chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

		#2.2)CHART LEGEND
		if has_legend:
			chart.has_legend = has_legend
			chart.legend.position = XL_LEGEND_POSITION.BOTTOM
			chart.legend.include_in_layout = False
			chart.legend.font.size = Pt(legend_font)

		#2.3)CÁC SỐ PHẦN TRĂM TRONG CHART
		if has_data_labels:
			chart.plots[0].has_data_labels = True
			data_labels = chart.plots[0].data_labels
			data_labels.font.size = Pt(6)
			data_labels.number_format = '0%'
			data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

	# COLUMN STACK CHART
	def create_column_stack_chart(self,categories,dict_data,left,top,width,height,
																legend_font=12,title=None,has_legend=True,has_data_labels=False,number_format=''):

		chart_type = XL_CHART_TYPE.COLUMN_STACKED
		self.__create_normal_chart(chart_type,categories,dict_data,left,top,width,height,
															legend_font=legend_font,title=title,has_legend=has_legend,has_data_labels=has_data_labels,number_format=number_format)
	
	#LINE CHART
	def create_line_chart(self,categories,dict_data,left,top,width,height,
													legend_font=12,title=None,has_legend=True,has_data_labels=False,number_format=''):

		chart_type = XL_CHART_TYPE.LINE
		self.__create_normal_chart(chart_type,categories,dict_data,left,top,width,height,
																legend_font=legend_font,title=title,has_legend=has_legend,has_data_labels=has_data_labels,number_format=number_format)

	#COLUMN CHART
	def create_column_chart(self,categories,dict_data,left,top,width,height,
													legend_font=12,title=None,has_legend=True,has_data_labels=False,number_format=''):

		chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
		self.__create_normal_chart(chart_type,categories,dict_data,left,top,width,height,
																legend_font=legend_font,title=title,has_legend=has_legend,
																has_data_labels=has_data_labels,number_format=number_format)

	#BAR CHART
	def create_bar_chart(self,categories,dict_data,left,top,width,height,
													legend_font=12,title=None,has_legend=True,has_data_labels=False,number_format=''):

		chart_type = XL_CHART_TYPE.BAR_CLUSTERED
		self.__create_normal_chart(chart_type,categories,dict_data,left,top,width,height,
																legend_font=legend_font,title=title,has_legend=has_legend,
																has_data_labels=has_data_labels,number_format=number_format)

	'''
	Hàm vẽ chart cấp 2 có thể vẽ chart và thêm data table ở phía dưới 
	'''
	# BAR CHART
	def create_bar_chart_with_data_table_follow_label(self,categories,dict_data,left,top,width,height,
																											legend_font=12,title=None,has_legend=True,
																											has_data_labels=False,number_format='',table_font_size=12,
																											table_columns_width={}):
		#1)Tạo chart
		self.create_bar_chart(categories,dict_data,left,top,width,height,number_format=number_format,
														legend_font=legend_font,title=title,has_legend=has_legend,has_data_labels=has_data_labels)

		#2)Vẽ table
		header = ['年月',title]
		cell_width = width
		cell_height = height/3
		table_left = left
		table_top = top + height
		self.__create_chart_table_follow_label(header,categories,dict_data,table_left,table_top,
																					cell_width,cell_height,font_size=table_font_size,
																					table_columns_width=table_columns_width)

	# COLUMN CHART
	def create_column_chart_with_data_table_follow_label(self,categories,dict_data,left,top,width,height,
																											legend_font=12,title=None,has_legend=True,
																											has_data_labels=False,number_format='',table_font_size=12,
																											table_columns_width={},dict_data_table = None):
		#1)Tạo chart
		self.create_column_chart(categories,dict_data,left,top,width,height,number_format=number_format,
														legend_font=legend_font,title=title,has_legend=has_legend,has_data_labels=has_data_labels)

		#2)Vẽ table
		header = ['年月',title]
		cell_width = width
		cell_height = height/3
		table_left = left
		table_top = top + height
		self.__create_chart_table_follow_label(header,categories,dict_data,table_left,table_top,
																					cell_width,cell_height,font_size=table_font_size,
																					table_columns_width=table_columns_width,dict_data_table = dict_data_table)

	def create_column_chart_with_data_table_follow_categories(self,categories,dict_data,left,top,width,height,
																														legend_font=12,title=None,has_legend=True,
																														has_data_labels=False,number_format='',table_font_size=12,
																														table_columns_width={}):
		#1)Tạo chart
		self.create_column_chart(categories,dict_data,left,top,width,height,number_format=number_format,
														legend_font=legend_font,title=title,has_legend=has_legend,has_data_labels=has_data_labels
														)

		#2)Vẽ table
		header = ['',title]
		cell_width = width
		cell_height = height/3
		table_left = left
		table_top = top + height
		self.__create_chart_table_follow_categories(header,categories,dict_data,table_left,table_top,
																								cell_width,cell_height,font_size=table_font_size,
																								table_columns_width=table_columns_width)

	# LINE CHART
	def create_line_chart_with_data_table_follow_label(self,categories,dict_data,left,top,width,height,
																										legend_font=12,title=None,has_legend=True,
																										has_data_labels=False,number_format='',table_font_size=12,
																										table_columns_width={}):
		#1)Tạo chart
		self.create_line_chart(categories,dict_data,left,top,width,height,number_format=number_format,
													legend_font=legend_font,title=title,has_legend=has_legend,has_data_labels=has_data_labels)

		#2)Vẽ table
		header = ['年月',title]
		cell_width = width
		cell_height = height/3
		table_left = left
		table_top = top + height
		self.__create_chart_table_follow_label(header,categories,dict_data,table_left,table_top,
																						cell_width,cell_height,font_size=table_font_size,
																						table_columns_width=table_columns_width)

	def create_line_chart_with_data_table_follow_categories(self,categories,dict_data,left,top,width,height,
																													legend_font=12,title=None,has_legend=True,
																													has_data_labels=False,number_format='',table_font_size=12,
																													table_columns_width={}):
		#1)Tạo chart
		self.create_line_chart(categories,dict_data,left,top,width,height,number_format=number_format,
														legend_font=legend_font,title=title,has_legend=has_legend,has_data_labels=has_data_labels)

		#2)Vẽ table
		header = ['',title]
		cell_width = width
		cell_height = height/3
		table_left = left
		table_top = top + height
		self.__create_chart_table_follow_categories(header,categories,dict_data,table_left,table_top,
																								cell_width,cell_height,font_size=table_font_size,
																								table_columns_width=table_columns_width)

	# COLUMN STACK CHART
	def create_column_stack_chart_with_data_table_follow_categories(self,categories,dict_data,left,top,width,height,
																																	legend_font=12,title=None,has_legend=True,
																																	has_data_labels=False,number_format='',table_font_size=12,
																																	table_columns_width={}):
		#1)Tạo chart
		self.create_column_stack_chart(categories,dict_data,left,top,width,height,number_format=number_format,
																			legend_font=legend_font,title=title,has_legend=has_legend,has_data_labels=has_data_labels)

		#2)Vẽ table
		header = ['',title]
		cell_width = width
		cell_height = height/3
		table_left = left
		table_top = top + height
		self.__create_chart_table_follow_categories(header,categories,dict_data,table_left,table_top,cell_width,
																								cell_height,font_size=table_font_size,table_columns_width=table_columns_width)

	def create_column_stack_chart_with_data_table_follow_both_label_and_categories(self,categories,dict_data,left,top,width,height,
																																	legend_font=12,title=None,has_legend=True,
																																	has_data_labels=False,number_format='',table_font_size=12,
																																	table_columns_width={}):
		#1)Tạo chart
		self.create_column_stack_chart(categories,dict_data,left,top,width,height,number_format=number_format,
																			legend_font=legend_font,title=title,has_legend=has_legend,has_data_labels=has_data_labels)

		#2)Vẽ table
		header = [title]
		header.extend(categories)
		cell_width = width
		cell_height = height/3
		table_left = left
		table_top = top + height
		self.create_chart_table_follow_both_label_and_categories(header,categories,dict_data,table_left,table_top,
																														cell_width,cell_height,font_size=table_font_size,
																														table_columns_width=table_columns_width)
	
	# PIE CHART
	def create_pie_chart_with_data_table_follow_label(self,header,categories,dict_data,left,top,width,height,
																										legend_font=12,title=None,has_legend=True,table_font_size=12,
																										has_table=True,
																										has_data_labels=False,
																										table_columns_width={}):
		#1)Tạo chart
		self.create_pie_chart(categories,dict_data,left,top,width,height,
													legend_font=legend_font,title=title,has_legend=has_legend,has_data_labels=has_data_labels)

		#2)Vẽ table  
		if has_table:
			# header = ['月',title]
			cell_width = width
			cell_height = height/3
			table_left = left
			table_top = top + height
			self.create_pie_chart_table(header,categories,dict_data,table_left,table_top,
																				cell_width,cell_height,font_size=table_font_size,table_columns_width=table_columns_width)
  
